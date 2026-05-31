#!/usr/bin/env python3
"""Extract all PPTX content to JSON files in the workspace."""
import sys
import os
import json
import subprocess

# First check if python-pptx is installed
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed, installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

workspace = "/root/.hermes/kanban/workspaces/t_96e95330"
base = "/root/work/dabt/dabt-tutor/reference/exam-materials/practice-exams/abt-2015-recert"

os.makedirs(workspace, exist_ok=True)

for fname in sorted(os.listdir(base)):
    if fname.endswith('.pptx') and not fname.endswith('.hms-bak'):
        full_path = os.path.join(base, fname)
        print(f"Processing: {fname}")
        prs = Presentation(full_path)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = p.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides_data.append({"slide": i+1, "texts": texts})
        outname = fname.rsplit('.', 1)[0] + '.json'
        outpath = os.path.join(workspace, outname)
        with open(outpath, 'w') as f:
            json.dump(slides_data, f, indent=2)
        print(f"  -> {outname} ({len(slides_data)} slides)")

print("\nDone! All JSON exports generated.")
