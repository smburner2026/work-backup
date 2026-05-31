#!/usr/bin/env python3
"""Extract text content from all 2015 Part B and Part C PPTX answer files."""
import sys, os, json, glob, re

SRC = "/root/work/dabt/dabt-tutor/reference/exam-materials/practice-exams/abt-2015-recert"

# Try to use python-pptx
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed. Trying to install...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

def extract_slides(pptx_path):
    """Extract text from each slide of a PPTX file."""
    prs = Presentation(pptx_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))
        slides_text.append({
            "slide_number": i + 1,
            "texts": slide_texts
        })
    return slides_text

def main():
    # Define the files to process
    files = {
        "part_b_q1_8": "Szabo questions 1_8.pptx",
        "part_b_q9_16": "2015 Part b JG_9-16.pptx",
        "part_b_q17_24": "Natalia_Recert part B_17-24.pptx",
        "part_b_q24_40": "2015 part b 24-40 ADL_AM.pptx",
        "part_c_q1_8": "DS_2015 part C questions 1_8.pptx",
        "part_c_q9_16": "JG_part C 9-16.pptx",
    }
    
    output = {}
    for key, fname in files.items():
        path = os.path.join(SRC, fname)
        if not os.path.exists(path):
            print(f"WARNING: {path} not found!")
            continue
        print(f"Processing {fname}...")
        output[key] = extract_slides(path)
    
    outpath = "/root/work/pptx_extracted.json"
    with open(outpath, "w") as f:
        json.dump(output, f, indent=2, ensure_ascii=False)
    print(f"\nExtraction complete! Saved to {outpath}")
    print(f"Files processed: {list(output.keys())}")

if __name__ == "__main__":
    main()
